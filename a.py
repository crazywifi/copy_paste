# Core functionality
import os
import re
import time
import json
import sys
import requests
import concurrent.futures
from datetime import datetime
from urllib.parse import urlparse, parse_qs

# GUI components
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, ttk, Menu

# Web automation
from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service

# System interactions
import subprocess
import platform
import webbrowser

# Document parsing (for search functionality)
import fitz  # PyMuPDF for PDFs
import docx
import pptx
import pandas as pd

# Constants
RESUME_FOLDER = "Resume_Download"
RESUME_URL = "https://apim.people.deloitte/personresumes?email="
LOGIN_URL = "https://people.deloitte/profile/"
EMAIL_URL_FILE = "Email_URL_Write.txt"
RESUME_URL_FILE = "ResumeDownloadURL.txt"
MAX_WORKERS = 20
TOKEN_REGEX = r'(\"secret\":\")(.*?)\"'

# Global variables
authorization_token = None

class TextRedirector:
    """Redirects console output to a tkinter text widget"""
    def __init__(self, text_widget):
        self.text_widget = text_widget

    def write(self, text):
        self.text_widget.insert(tk.END, text, ("color",))
        self.text_widget.see(tk.END)  # Scroll to the end of the text
        
    def flush(self):
        pass  # Required to satisfy the sys.stdout interface


class ResumeDownloader:
    """Main application class for the Deloitte Resume Downloader"""
    
    def __init__(self, root):
        self.root = root
        self.email_id_file = None
        self.setup_ui()
        self.create_folder()
        
    def setup_ui(self):
        """Set up the user interface"""
        self.root.title("Deloitte Resume Downloader v1.1")
        self.root.geometry("400x500")
        self.root.configure(bg="#ddc8cf")
        
        try:
            self.root.iconbitmap(r".\icon.ico")
        except tk.TclError:
            print("Icon file not found or not a valid icon file.")
        
        # File selection
        tk.Label(self.root, text="Select Email txt File:", bg="#ddc8cf").pack(pady=10)
        self.email_id_entry = tk.Entry(self.root, width=40)
        self.email_id_entry.pack(pady=5)
        
        tk.Label(self.root, text="Follow steps one by one", bg="#ddc8cf").pack(pady=5)
        
        # Buttons
        tk.Button(self.root, text="Browse", command=self.browse_file).pack(pady=5)
        tk.Button(self.root, text="Sign In", command=self.login_to_microsoft).pack(pady=5)
        tk.Button(self.root, text="Download", command=self.download_resume).pack(pady=5)
        tk.Button(self.root, text="Search", command=self.search_resume).pack(pady=5)
        tk.Button(self.root, text="Close", command=self.on_closing).pack(pady=5)
        
        # Footer information
        tk.Label(self.root, text="Created By Rishabh Sharma", fg="red", bg="#ddc8cf").pack(pady=2)
        tk.Label(self.root, text="For any feedback:", fg="red", font=("Arial", 10), bg="#ddc8cf").pack(pady=1)
        
        email_label = tk.Label(self.root, text="rishabhsharma96@deloitte.com", fg="red", 
                            cursor="hand2", font=("Arial", 10, "underline"), bg="#ddc8cf")
        email_label.pack(pady=2)
        email_label.bind("<Button-1>", lambda e: webbrowser.open("mailto:rishabhsharma96@deloitte.com"))
        
        # Progress frame
        self.progress_frame = tk.Frame(self.root, bg="#ddc8cf")
        self.progress_frame.pack(pady=5, fill=tk.X, padx=10)
        
        self.progress_label = tk.Label(self.progress_frame, text="Download Progress:", bg="#ddc8cf")
        self.progress_label.pack(side=tk.TOP, anchor="w")
        
        self.progress_bar = ttk.Progressbar(self.progress_frame, orient="horizontal", length=380, mode="determinate")
        self.progress_bar.pack(fill=tk.X)
        
        # Console output
        self.output_text = tk.Text(self.root, wrap="word", height=15)
        self.output_text.pack(fill=tk.BOTH, expand=True)
        self.output_text.tag_configure("color", foreground="green")
        sys.stdout = TextRedirector(self.output_text)
        
    def browse_file(self):
        """Open file browser to select email ID file"""
        file_path = filedialog.askopenfilename(filetypes=[("Text files", "*.txt")])
        if file_path:
            self.email_id_entry.delete(0, tk.END)
            self.email_id_entry.insert(tk.END, file_path)
            self.email_id_file = self.email_id_entry.get()
    
    def login_to_microsoft(self):
        """Handle Microsoft login process and extract authorization token"""
        global authorization_token
        self.create_folder()
        
        print("Starting login process...")
        
        # Get paths for Chrome executables
        current_directory = os.path.dirname(os.path.abspath(__file__))
        chrome_executable_path = os.path.join(current_directory, 'chrome-win', 'chrome.exe')
        chromedriver_executable_path = os.path.join(current_directory, 'chrome', 'chromedriver.exe')
        os.environ['PATH'] += os.pathsep + chromedriver_executable_path
        
        # Configure Chrome options
        chrome_options = webdriver.ChromeOptions()
        chrome_options.binary_location = chrome_executable_path
        
        # Add performance options
        chrome_options.add_argument('--disable-extensions')
        chrome_options.add_argument('--disable-gpu')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument('--no-sandbox')
        
        driver = None
        try:
            driver = webdriver.Chrome(options=chrome_options)
            driver.get(LOGIN_URL)
            
            # Wait for authorization token
            token_found = False
            max_attempts = 60  # 2 minutes timeout (2sec * 60)
            attempts = 0
            
            while not token_found and attempts < max_attempts:
                session_storage_data = driver.execute_script("""
                    var sessionStorageData = {};
                    for (var i = 0; i < sessionStorage.length; i++) {
                        var key = sessionStorage.key(i);
                        var value = sessionStorage.getItem(key);
                        sessionStorageData[key] = value;
                    }
                    return sessionStorageData;
                """)
                
                match = re.search(TOKEN_REGEX, str(session_storage_data), re.MULTILINE)
                if match:
                    authorization_token = match.group(2)
                    token_found = True
                    print("Login Successful!")
                    print(f"Authorization Token: {authorization_token[:10]}...")
                    break
                
                attempts += 1
                time.sleep(2)
                
            if not token_found:
                print("Login timed out. Please try again.")
                
        except Exception as e:
            print(f"Login error: {e}")
        finally:
            if driver:
                driver.quit()
    
    def download_resume(self):
        """Download resumes for the provided email IDs"""
        global authorization_token
        
        if not authorization_token:
            print("Error: Not logged in. Please sign in first.")
            return
            
        if not self.email_id_file:
            print("Error: No email file selected. Please browse and select a file.")
            return
            
        print("\nInitiating Download Process...")
        
        headers = {"Authorization": f"Bearer {authorization_token}"}
        
        try:
            # Step 1: Read email IDs and create URLs
            print("Reading Emails From EmailId.txt")
            with open(self.email_id_file, 'r') as email_file:
                email_ids = [email.strip() for email in email_file.readlines()]
            
            # Setup progress tracking for UI
            total_emails = len(email_ids)
            self.progress_bar["maximum"] = total_emails
            self.progress_bar["value"] = 0
            self.progress_label.config(text="Preparing Email URLs:")
            self.root.update_idletasks()
            
            # Create URLs file with deduplication
            self.prepare_email_urls(email_ids)
            
            # Step 2: Fetch resume metadata
            print("Fetching resume information...")
            urls = self.read_urls_from_file(EMAIL_URL_FILE)
            self.fetch_resume_metadata(urls, headers)
            
            # Step 3: Download actual resume files
            print("Downloading resume files...")
            resume_urls = self.read_urls_from_file(RESUME_URL_FILE)
            
            # Reset progress bar for downloads
            total_downloads = len(resume_urls)
            self.progress_bar["maximum"] = total_downloads
            self.progress_bar["value"] = 0
            self.progress_label.config(text="Downloading Resumes:")
            self.root.update_idletasks()
            
            self.download_resume_files(resume_urls, headers)
            
            # Cleanup
            self.cleanup_temp_files()
            print("Download process completed successfully!")
            self.progress_label.config(text="Download Complete!")
            
        except Exception as e:
            print(f"Error during download process: {e}")
            self.cleanup_temp_files()
    
    def prepare_email_urls(self, email_ids):
        """Prepare email URLs file with deduplication"""
        try:
            # Read existing URLs if file exists
            existing_emails = set()
            if os.path.exists(EMAIL_URL_FILE):
                with open(EMAIL_URL_FILE, 'r') as url_file:
                    existing_emails = set(url_file.read().splitlines())
            
            # Add new URLs with deduplication
            with open(EMAIL_URL_FILE, 'a') as url_file:
                for i, email in enumerate(email_ids):
                    email_url = f"{RESUME_URL}{email}"
                    if email_url not in existing_emails:
                        url_file.write(f"{email_url}\n")
                        existing_emails.add(email_url)
                    
                    # Update progress
                    self.progress_bar["value"] = i + 1
                    self.root.update_idletasks()
            
        except Exception as e:
            print(f"Error preparing email URLs: {e}")
    
    def fetch_resume_metadata(self, urls, headers):
        """Fetch resume metadata using concurrent requests"""
        
        # Reset progress
        self.progress_bar["maximum"] = len(urls)
        self.progress_bar["value"] = 0
        self.progress_label.config(text="Fetching Resume Info:")
        self.root.update_idletasks()
        
        # Create or clear the resume URL file
        open(RESUME_URL_FILE, 'w').close()
        
        completed = 0
        
        def process_response(future):
            nonlocal completed
            url = future_to_url[future]
            try:
                response = future.result()
                json_data = response.json()
                
                # Find the most recent resume
                if "data" in json_data and json_data["data"]:
                    latest_resume = self.find_latest_resume(json_data["data"])
                    
                    if latest_resume:
                        document_url = latest_resume["documentUrl"].strip()
                        print(f"Found latest resume: {latest_resume['fileName']} created on {latest_resume['creationDate']}")
                        
                        with open(RESUME_URL_FILE, 'a') as url_file:
                            url_file.write(f"{document_url}\n")
                
            except Exception as e:
                print(f"Error processing {url}: {e}")
            finally:
                # Update progress bar
                completed += 1
                self.progress_bar["value"] = completed
                self.root.update_idletasks()
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            future_to_url = {executor.submit(self.send_get_request, url, headers): url for url in urls}
            
            # Process results as they complete
            for future in concurrent.futures.as_completed(future_to_url):
                process_response(future)
    
    def download_resume_files(self, urls, headers):
        """Download resume files using concurrent requests"""
        completed = 0
        
        def process_download(future):
            nonlocal completed
            url = future_to_url[future]
            try:
                response = future.result()
                content_disposition = response.headers.get("content-disposition", "")
                
                # Extract filename
                filename_match = re.search(r'filename="([^"]+)"', content_disposition)
                if filename_match:
                    filename = filename_match.group(1)
                    print(f"Downloading: {filename}")
                    
                    # Save file
                    save_path = os.path.join(RESUME_FOLDER, filename)
                    with open(save_path, 'wb') as file:
                        file.write(response.content)
                        print(f"Saved to {save_path}")
                
            except Exception as e:
                print(f"Error downloading from {url}: {e}")
            finally:
                # Update progress
                completed += 1
                self.progress_bar["value"] = completed
                self.root.update_idletasks()
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            future_to_url = {executor.submit(self.download_file, url, headers): url for url in urls}
            
            # Process results as they complete
            for future in concurrent.futures.as_completed(future_to_url):
                process_download(future)
    
    def search_resume(self):
        """Open the search window"""
        search_window = tk.Toplevel(self.root)
        search_window.title("Resume Search")
        search_window.geometry("500x400")
        search_window.configure(bg="#ddc8cf")
        
        try:
            search_window.iconbitmap(r".\icon1.ico")
        except tk.TclError:
            print("Icon file 'icon1.ico' not found or not a valid icon file.")
        
        # Search Entry Field
        tk.Label(search_window, text="Enter search keywords:", bg="#ddc8cf").pack(pady=5)
        search_entry = tk.Entry(search_window, width=50)
        search_entry.pack(pady=5)
        
        # Operator Buttons
        button_frame = tk.Frame(search_window, bg="#ddc8cf")
        button_frame.pack(pady=5)
        
        tk.Button(button_frame, text="AND", 
                command=lambda: self.append_operator("AND", search_entry)).pack(side=tk.LEFT, padx=5)
        tk.Button(button_frame, text="OR", 
                command=lambda: self.append_operator("OR", search_entry)).pack(side=tk.LEFT, padx=5)
        tk.Button(button_frame, text="NOT", 
                command=lambda: self.append_operator("NOT", search_entry)).pack(side=tk.LEFT, padx=5)
        
        # Folder Display
        self.folder_label = tk.Label(search_window, text=f"Folder: {RESUME_FOLDER}", bg="#ddc8cf")
        self.folder_label.pack(pady=5)
        
        # Change Folder Button
        tk.Button(search_window, text="Change Folder", 
                command=self.browse_folder).pack(pady=5)
        
        # Search Button
        tk.Button(search_window, text="Search", 
                command=lambda: self.perform_search(result_text, search_entry.get())).pack(pady=5)
        
        # Result Text Area
        result_text = scrolledtext.ScrolledText(search_window, height=10, width=60, cursor="arrow")
        result_text.pack(pady=5)
        
        # Bind events
        result_text.bind("<Motion>", lambda event: self.on_hover(event, result_text))
        result_text.bind("<Double-Button-1>", lambda event: self.open_file(event, result_text))
    
    # Helper methods
    def create_folder(self):
        """Create resume download folder if it doesn't exist"""
        if not os.path.exists(RESUME_FOLDER):
            os.mkdir(RESUME_FOLDER)
            print(f"Folder {RESUME_FOLDER} created!")
    
    def on_closing(self):
        """Handle application closing"""
        self.cleanup_temp_files()
        if messagebox.askokcancel("Quit", "Do you want to quit?"):
            self.root.destroy()
    
    def cleanup_temp_files(self):
        """Remove temporary files"""
        for filename in [EMAIL_URL_FILE, RESUME_URL_FILE]:
            if os.path.exists(filename):
                try:
                    os.remove(filename)
                except Exception as e:
                    print(f"Error removing {filename}: {e}")
    
    def read_urls_from_file(self, file_path):
        """Read URLs from a file"""
        if not os.path.exists(file_path):
            return []
            
        with open(file_path, "r") as file:
            return [line.strip() for line in file if line.strip()]
    
    def send_get_request(self, url, headers):
        """Send GET request with error handling and retries"""
        max_retries = 3
        for attempt in range(max_retries):
            try:
                response = requests.get(url, headers=headers, verify=False, timeout=30)
                return response
            except requests.exceptions.RequestException as e:
                if attempt < max_retries - 1:
                    print(f"Retry {attempt+1}/{max_retries} for {url}: {e}")
                    time.sleep(2)
                else:
                    raise
    
    def download_file(self, url, headers):
        """Download file with error handling and retries"""
        max_retries = 3
        for attempt in range(max_retries):
            try:
                response = requests.get(url, headers=headers, timeout=30)
                return response
            except requests.exceptions.RequestException as e:
                if attempt < max_retries - 1:
                    print(f"Retry {attempt+1}/{max_retries} for {url}: {e}")
                    time.sleep(2)
                else:
                    raise
    
    def find_latest_resume(self, resume_data):
        """Find the most recent resume from the data"""
        latest_resume = None
        latest_date = None
        
        for item in resume_data:
            creation_date = item.get("creationDate")
            if creation_date:
                try:
                    current_date = datetime.strptime(creation_date, "%Y-%m-%dT%H:%M:%S")
                    if latest_date is None or current_date > latest_date:
                        latest_date = current_date
                        latest_resume = item
                except ValueError:
                    pass  # Skip items with invalid date format
        
        return latest_resume
    
    # Search functionality methods
    def perform_search(self, result_text, query):
        """Perform search across resume files"""
        if not query.strip():
            messagebox.showerror("Error", "Please enter a search query.")
            return
        
        formatted_query = self.format_search_query(query)
        print(f"Searching for: {formatted_query}")
        
        matching_files = []
        total_files = len([f for f in os.listdir(RESUME_FOLDER) if os.path.isfile(os.path.join(RESUME_FOLDER, f))])
        
        # Show a loading indicator
        result_text.delete("1.0", tk.END)
        result_text.insert(tk.END, f"Searching {total_files} files...\n")
        self.root.update_idletasks()
        
        for file in os.listdir(RESUME_FOLDER):
            filepath = os.path.join(RESUME_FOLDER, file)
            if os.path.isfile(filepath):
                # Update the search status
                result_text.delete("1.0", "2.0")
                result_text.insert("1.0", f"Searching {file}...\n")
                self.root.update_idletasks()
                
                try:
                    text = self.extract_text(filepath)
                    if self.boolean_search(text, formatted_query):
                        matching_files.append(file)
                except Exception as e:
                    print(f"Error searching {file}: {e}")
        
        # Display results
        result_text.delete("1.0", tk.END)
        if matching_files:
            result_text.insert(tk.END, f"Found {len(matching_files)} matches:\n\n")
            result_text.insert(tk.END, "\n".join(matching_files))
        else:
            result_text.insert(tk.END, "No matching resumes found.")
    
    def extract_text(self, filepath):
        """Extract text from various document formats"""
        ext = filepath.lower().split(".")[-1]
        
        if ext == "pdf":
            return self.extract_text_from_pdf(filepath)
        elif ext == "docx":
            return self.extract_text_from_docx(filepath)
        elif ext in ["pptx", "ppt"]:
            return self.extract_text_from_pptx(filepath)
        elif ext in ["xls", "xlsx"]:
            return self.extract_text_from_excel(filepath)
        else:
            print(f"Unsupported format: {filepath}")
            return ""
    
    def extract_text_from_pdf(self, filepath):
        """Extract text from PDF file"""
        text = ""
        try:
            with fitz.open(filepath) as doc:
                for page in doc:
                    text += page.get_text("text") + "\n"
        except Exception as e:
            print(f"Error reading PDF {filepath}: {e}")
        return text
    
    def extract_text_from_docx(self, filepath):
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            print(f"Error reading DOCX {filepath}: {e}")
            return ""
    
    def extract_text_from_pptx(self, filepath):
        """Extract text from PPTX file"""
        try:
            presentation = pptx.Presentation(filepath)
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error reading PPTX {filepath}: {e}")
            return ""
    
    def extract_text_from_excel(self, filepath):
        """Extract text from Excel file"""
        try:
            df = pd.read_excel(filepath, sheet_name=None)
            text = []
            for sheet_name, sheet_df in df.items():
                text.append(f"Sheet: {sheet_name}")
                text.append(sheet_df.to_string())
            return "\n".join(text)
        except Exception as e:
            print(f"Error reading Excel {filepath}: {e}")
            return ""
    
    def format_search_query(self, query):
        """Format search query with appropriate operators"""
        words = query.split()
        operators = {"AND", "OR", "NOT"}
        formatted_query = []
        temp_phrase = []
        
        for word in words:
            if word.upper() in operators:
                if temp_phrase:
                    formatted_query.append(f'"{" ".join(temp_phrase)}"')
                    temp_phrase = []
                formatted_query.append(word)
            else:
                temp_phrase.append(word)
        
        if temp_phrase:
            formatted_query.append(f'"{" ".join(temp_phrase)}"')
        
        final_query = " ".join(formatted_query)
        
        # If no logical operators, wrap entire query in quotes
        if not any(op in final_query for op in operators) and " " in final_query:
            if not (final_query.startswith('"') and final_query.endswith('"')):
                final_query = f'"{final_query}"'
        
        return final_query
    
    def boolean_search(self, text, query):
        """Perform boolean search on text"""
        query = query.replace("AND", "&").replace("OR", "|").replace("NOT", "~")
        
        words = re.findall(r'"[^"]+"|\w+', query)
        
        for word in words:
            word_cleaned = word.strip('"')
            if word_cleaned.lower() not in text.lower():
                query = query.replace(word, "False")
            else:
                query = query.replace(word, "True")
        
        try:
            return eval(query)
        except:
            return False
    
    def append_operator(self, op, search_entry):
        """Insert a logical operator into the search entry field"""
        text = search_entry.get().strip()
        
        if text and not text.endswith(("AND", "OR", "NOT")):
            search_entry.insert(tk.END, f" {op} ")
        elif not text:
            messagebox.showerror("Error", "Please enter a search keyword before adding an operator.")
    
    def browse_folder(self):
        """Browse for a different folder to search in"""
        global RESUME_FOLDER
        folder_selected = filedialog.askdirectory()
        if folder_selected:
            RESUME_FOLDER = folder_selected
            self.folder_label.config(text=f"Folder: {RESUME_FOLDER}")
    
    def on_hover(self, event, text_widget):
        """Handle mouse hover over result text"""
        try:
            text_widget.tag_remove("hover", "1.0", tk.END)
            index = text_widget.index(f"@{event.x},{event.y}")
            line_start = f"{index.split('.')[0]}.0"
            line_end = f"{index.split('.')[0]}.end"
            
            text_widget.tag_add("hover", line_start, line_end)
            text_widget.tag_config("hover", background="lightblue")
        except Exception as e:
            print(f"Error in hover effect: {e}")
    
    def open_file(self, event, text_widget):
        """Open the selected file"""
        try:
            index = text_widget.index(tk.CURRENT)
            line_start = f"{index.split('.')[0]}.0"
            line_end = f"{index.split('.')[0]}.end"
            selected_text = text_widget.get(line_start, line_end).strip()
            
            if selected_text:
                filepath = os.path.join(RESUME_FOLDER, selected_text)
                if os.path.exists(filepath):
                    if platform.system() == "Windows":
                        os.startfile(filepath)
                    elif platform.system() == "Darwin":  # macOS
                        subprocess.run(["open", filepath])
                    else:  # Linux
                        subprocess.run(["xdg-open", filepath])
        except Exception as e:
            print(f"Error opening file: {e}")


# Main application entry point
if __name__ == "__main__":
    root = tk.Tk()
    app = ResumeDownloader(root)
    root.protocol("WM_DELETE_WINDOW", app.on_closing)
    root.mainloop()
